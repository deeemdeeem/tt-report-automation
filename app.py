import os
import io
import tempfile
from datetime import datetime
from copy import deepcopy

from flask import Flask, render_template_string, request, send_file, redirect, url_for, flash
from werkzeug.utils import secure_filename

import pandas as pd
import numbers
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_TICK_MARK
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor



# ---------------- Flask setup ----------------
app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "dev-key")

# PPT Template on file directory
PPT_TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "TT_report.pptx")
ALLOWED_EXCEL_EXTS = {".xlsm", ".xlsx"}
SHAREPOINT_TEMPLATE_URL = "https://aro36579709.sharepoint.com/:x:/s/ABResources/IQCj5Ih6CsWDSL1vFuuGeyy5AYo0yCTvekxuIjFJrEkA0EA?e=RYNca1"

HTML = """
<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <meta name="viewport" content="width=device-width, initial-scale=1" />
  <title>TT Report Generator</title>
  <link rel="preconnect" href="https://fonts.googleapis.com">
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700&display=swap" rel="stylesheet">
  <style>
    :root{
      --bg:#0b1220;        /* deep navy */
      --card:#0f172a;      /* slate-900 */
      --muted:#94a3b8;     /* slate-400 */
      --text:#e2e8f0;      /* slate-200 */
      --accent:#6d28d9;    /* violet-700 */
      --accent-hover:#7c3aed;
      --border:#1f2937;    /* slate-800 */
      --good:#10b981;      /* emerald */
      --warn:#f59e0b;      /* amber */
    }
    *{ box-sizing: border-box; }
    html,body{ height:100%; }
    body{
      margin:0; background:var(--bg); color:var(--text);
      font-family: Inter, system-ui, -apple-system, Segoe UI, Roboto, Arial, sans-serif;
    }
    .header{
      position: sticky; top:0; background:rgba(15,23,42,.8);
      backdrop-filter: blur(8px);
      border-bottom:1px solid var(--border);
    }
    .header-inner{
      max-width:900px; margin:0 auto; padding:12px 20px; display:flex; align-items:center; gap:12px;
    }
    .header img{ height:28px; display:block; }
    .wrap{ max-width:900px; margin:36px auto; padding:0 20px; }
    .card{
      background:var(--card); border:1px solid var(--border);
      border-radius:16px; padding:24px; box-shadow:0 10px 30px rgba(0,0,0,.35);
    }
    h1, h2{ margin:0 0 6px; font-weight:700; letter-spacing:.2px; }
    h1{ font-size:26px; }
    h2{ font-size:20px; color:#fff; }
    p.sub{ margin:0 0 18px; color:var(--muted); }
    form{ display:flex; flex-wrap:wrap;ap:12px; align-items:center; }
    .file-wrap{
      position:relative; display:flex; align-items:center; gap:10px;
      border:1px dashed var(--border); background:rgba(2,6,23,.5);
      padding:12px 14px; border-radius:12px; min-width:320px; flex:1;
    }
    input[type=file]{ position:absolute; inset:0; opacity:0; cursor:pointer; }
    .file-label{ color:var(--muted); font-size:14px; }
    .file-name{ font-size:14px; color:var(--text); white-space:nowrap; overflow:hidden; text-overflow:ellipsis; }
    .btn{
      border:none; padding:12px 16px; border-radius:12px; font-weight:600; cursor:pointer;
      transition: transform .04s ease;
    }
    .btn:active{ transform: translateY(1px); }
    .btn-primary{ background:var(--accent); color:#fff; }
    .btn-primary:hover{ background:var(--accent-hover); }
    .btn-ghost{ background:transparent; border:1px solid var(--border); color:var(--text); }
    .hint{ margin-top:10px; color:var(--muted); font-size:12px; }
    .flash{
      background:rgba(245,158,11,.1); color:#fde68a;
      border:1px solid rgba(245,158,11,.35);
      padding:10px 12px; border-radius:10px; margin-bottom:12px;
    }
    /* --- Loading overlay + spinner --- */
    .overlay{
      position: fixed; inset:0; display:none; place-items:center;
      background: rgba(2,6,23,.55); z-index: 9999;
    }
    .spinner{
      width:56px; height:56px; border-radius:50%;
      border:4px solid rgba(255,255,255,.15);
      border-top-color: #a78bfa; /* lighter violet */
      animation: spin 0.9s linear infinite;
      box-shadow: 0 0 30px rgba(167,139,250,.45);
    }
    @keyframes spin { to { transform: rotate(360deg); } }
  </style>
</head>
<body>
  <!-- Header with TruTrade logo -->
  <div class="header">
    <div class="header-inner">
      <img alt="TruTrade" src="https://trutradebeta.alexanderbabbage.com/static/media/trutrade_logo.8ab17e85dea03ec4e762.png" />
    </div>
  </div>

  <div class="wrap">
    <!-- Step 1: Download template -->
    <div class="card" style="margin-bottom:16px;">
      <h2>Step 1: Download Template</h2>
      <p class="sub">
        Access the <strong>TruTrade Executive Report Worksheet (.xlsm)</strong> from SharePoint. Paste your TruTrade datasets into the template and the built-in macros will automatically analyze each tab.<br>
        <em>Note:</em> after you load each dataset, allow a short processing time while the template finishes its analysis.
      </p>
      <a class="btn btn-primary" href="{{ url_for('download_template') }}" target="_blank" rel="noopener">Open Template on SharePoint</a>
    </div>

    <!-- Step 2: Upload & Generate -->
    <div class="card">
      <h2>Step 2: Upload & Generate Report</h2>
      <p class="sub">
        Upload the <strong>filled worksheet (.xlsm or .xlsx)</strong> here, then click <strong>Generate</strong>. We’ll populate the PowerPoint template and download a new <strong>.pptx</strong> for you.
        <br><em>Heads up:</em> maps and charts are <strong>not</strong> auto-generated in this version.
      </p>

      {% with messages = get_flashed_messages() %}
        {% if messages %}
          {% for m in messages %}<div class="flash">{{ m }}</div>{% endfor %}
        {% endif %}
      {% endwith %}

      <form id="genForm" action="{{ url_for('generate') }}" method="post" enctype="multipart/form-data" onreset="resetName();">
        <div class="file-wrap">
          <span class="file-label">Choose File</span>
          <span id="fileName" class="file-name">No file selected</span>
          <input id="fileInput" type="file" name="xlsm" accept=".xlsm,.xlsx" required>
        </div>

        <button id="genBtn" class="btn btn-primary" type="submit">Generate</button>
        <button id="clearBtn" class="btn btn-ghost" type="reset">Clear</button>
      </form>

      <div class="hint">Template in use: <code>{{ template_name }}</code></div>
    </div>
  </div>

  <!-- Loading overlay -->
  <div id="overlay" class="overlay" aria-hidden="true">
    <div class="spinner" role="status" aria-label="Generating report..."></div>
  </div>

  <script>
    const input   = document.getElementById('fileInput');
    const nameEl  = document.getElementById('fileName');
    const form    = document.getElementById('genForm');
    const genBtn  = document.getElementById('genBtn');
    const clearBtn= document.getElementById('clearBtn');
    const overlay = document.getElementById('overlay');

    function resetName(){ nameEl.textContent = 'No file selected'; }
    input.addEventListener('change', () => {
      nameEl.textContent = input.files.length ? input.files[0].name : 'No file selected';
    });

    function showLoading(){
      overlay.style.display = 'grid';
      genBtn.disabled = true;
      clearBtn.disabled = true;
      genBtn.textContent = 'Generating…';
      genBtn.setAttribute('aria-busy','true');
    }
    function hideLoading(){
      overlay.style.display = 'none';
      genBtn.disabled = false;
      clearBtn.disabled = false;
      genBtn.textContent = 'Generate';
      genBtn.removeAttribute('aria-busy');
    }

    // Intercept submit so we can show spinner and control the download
    form.addEventListener('submit', async (e) => {
      e.preventDefault();
      if (!input.files.length) return;

      try{
        showLoading();
        const fd = new FormData(form);
        const res = await fetch(form.action, { method: 'POST', body: fd });
        if(!res.ok){
          hideLoading();
          alert('Error generating report. Please check your file and try again.');
          return;
        }
        const blob = await res.blob();
        const url = URL.createObjectURL(blob);
        const a = document.createElement('a');
        const ts = new Date().toISOString().slice(0,16).replace('T','_');
        a.href = url;
        a.download = `TT_report_${ts}.pptx`;
        document.body.appendChild(a);
        a.click();
        a.remove();
        URL.revokeObjectURL(url);
      }catch(err){
        console.error(err);
        alert('Unexpected error. Please try again.');
      }finally{
        hideLoading();
      }
    });
  </script>
</body>
</html>

"""
# Chart Replacement Data Logic
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_TICK_MARK
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd

def update_mileage_chart(slide, df_mileage):
    """
    Horizontal bar chart in middle of slide using MileageDemo!Q4:Q26.
    - No category labels
    - No axis labels
    - Data labels at end of each bar, black, size 9
    - Bar color: red < 100, green > 100, grey = 100
    """

    # ---- 1) Slice the data from the DataFrame ----
    start_row = 2      # Excel row 4 (0-based)
    end_row   = 26     # Excel row 26 (exclusive)
    col_q_index = 16   # column Q

    cat_series = df_mileage.iloc[start_row:end_row, 0]
    val_series = df_mileage.iloc[start_row:end_row, col_q_index]

    # We don't care about category text; keep them blank
    categories = ["" for _ in range(len(val_series))]

    values = []
    for v in val_series:
        if pd.isna(v):
            values.append(None)
        else:
            try:
                values.append(round(float(v)))   # whole-number index
            except (TypeError, ValueError):
                values.append(None)

    # ---- 2) Build chart data ----
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Index", values)

    # ---- 3) Position & size: make it taller ----
    left = Inches(11.3)
    top = Inches(2.392)
    width = Inches(1.7)
    height = Inches(5.735)   

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        left, top, width, height,
        chart_data
    )
    chart = graphic_frame.chart

    # ---- Remove value axis line & labels ----
    chart.has_title = False
    val_axis = chart.value_axis
    val_axis.major_tick_mark = XL_TICK_MARK.NONE
    val_axis.minor_tick_mark = XL_TICK_MARK.NONE
    val_axis.tick_labels.number_format_is_linked = False
    val_axis.tick_labels.number_format = ";;;"  # hide numbers
    # hide the axis line itself
    val_axis.format.line.fill.background()
    # hide gridlines if present
    if val_axis.has_major_gridlines:
        val_axis.major_gridlines.format.line.fill.background()

    # ---- Remove category axis line & labels ----
    cat_axis = chart.category_axis
    cat_axis.major_tick_mark = XL_TICK_MARK.NONE
    cat_axis.minor_tick_mark = XL_TICK_MARK.NONE
    cat_axis.tick_labels.number_format_is_linked = False
    cat_axis.tick_labels.number_format = ";;;"
    cat_axis.format.line.fill.background()

    # ---- 5) Data labels at end of each bar ----
    plot = chart.plots[0]
    series = plot.series[0]

    plot.has_data_labels = True
    data_labels = series.data_labels
    data_labels.show_value = True
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END  # right end of bar

    # ---- 6) Color bars & style labels ----
    for idx, point in enumerate(series.points):
        v = values[idx]
        if v is None:
            continue

        lbl = point.data_label
        lbl.number_format = "0"         # no decimals
        lbl.font.size = Pt(9)           # font size 9
        lbl.font.color.rgb = RGBColor(0, 0, 0)  # black

        fill = point.format.fill
        fill.solid()

        if v < 100:
            fill.fore_color.rgb = RGBColor(220, 38, 38)   # red
        elif v > 100:
            fill.fore_color.rgb = RGBColor(22, 163, 74)   # green
        else:
            fill.fore_color.rgb = RGBColor(148, 163, 184) # grey

    chart.has_legend = False


def update_ethnicities_chart(slide, df_drawdemos):
    """
    Horizontal bar chart for ETHNIC MAKEUP using DrawDemo!W24:X28.
    - Categories from column W, values from column X (expected decimals rendered as %)
    - Dark-blue bars, value labels outside end
    """
    start_row = 22  # Excel row 24 (0-based)
    end_row = 27    # Excel row 28 (exclusive)
    cat_col = 22    # column W
    val_col = 23    # column X

    categories = []
    values = []
    for cat, val in zip(
        df_drawdemos.iloc[start_row:end_row, cat_col],
        df_drawdemos.iloc[start_row:end_row, val_col],
    ):
        categories.append("" if pd.isna(cat) else str(cat))
        try:
            values.append(None if pd.isna(val) else float(val))
        except Exception:
            values.append(None)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Ethnic Makeup", values)

    # Placement: lower-left area of slide 8 (per reference positioning)
    left = Inches(0.14)
    top = Inches(6.19)
    width = Inches(3.8)   # slightly wider bars
    height = Inches(2.1)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        left, top, width, height,
        chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = False
    chart.has_title = False

    # Minimal axes
    val_axis = chart.value_axis
    val_axis.major_tick_mark = XL_TICK_MARK.NONE
    val_axis.minor_tick_mark = XL_TICK_MARK.NONE
    val_axis.format.line.fill.background()
    val_axis.tick_labels.number_format_is_linked = False
    val_axis.tick_labels.number_format = ";;;"
    if val_axis.has_major_gridlines:
        val_axis.major_gridlines.format.line.fill.background()

    cat_axis = chart.category_axis
    cat_axis.major_tick_mark = XL_TICK_MARK.NONE
    cat_axis.minor_tick_mark = XL_TICK_MARK.NONE
    cat_axis.format.line.fill.background()
    cat_axis.tick_labels.font.size = Pt(11)
    cat_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)

    plot = chart.plots[0]
    series = plot.series[0]
    series.gap_width = 50  # tighter spacing between bars
    plot.has_data_labels = True
    data_labels = series.data_labels
    data_labels.show_value = True
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = "0%"
    data_labels.font.name = "Roboto"
    data_labels.font.size = Pt(11)
    data_labels.font.color.rgb = RGBColor(0, 0, 0)

    dark_blue = RGBColor(21, 57, 96)
    for point in series.points:
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = dark_blue


def update_frequency_chart(slide, df_frequency):
    """
    Horizontal bar chart for Frequency (Frequency!J22:K31) placed on slide 19.
    - Categories from column J, values from column K (one decimal)
    - Dark blue bars, labels outside end
    """
    start_row =1  # Excel row 22 (0-based)
    end_row = 12    # Excel row 31 (exclusive)
    cat_col = 13     # column J
    val_col = 14    # column K

    categories = []
    values = []
    for cat, val in zip(
        df_frequency.iloc[start_row:end_row, cat_col],
        df_frequency.iloc[start_row:end_row, val_col],
    ):
        categories.append("" if pd.isna(cat) else str(cat))
        try:
            values.append(None if pd.isna(val) else float(val))
        except Exception:
            values.append(None)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Frequency", values)

    left = Inches(7)
    top = Inches(2.31)
    width = Inches(5.4)
    height = Inches(5.3)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        left, top, width, height,
        chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = False
    chart.has_title = False

    val_axis = chart.value_axis
    val_axis.major_tick_mark = XL_TICK_MARK.NONE
    val_axis.minor_tick_mark = XL_TICK_MARK.NONE
    val_axis.format.line.fill.background()
    val_axis.tick_labels.number_format_is_linked = False
    val_axis.tick_labels.number_format = ";;;"
    if val_axis.has_major_gridlines:
        val_axis.major_gridlines.format.line.fill.background()

    max_val = max((v for v in values if v is not None), default=0)
    val_axis.minimum_scale = 0
    val_axis.maximum_scale = max_val * 1.1 if max_val > 0 else 1

    cat_axis = chart.category_axis
    cat_axis.major_tick_mark = XL_TICK_MARK.NONE
    cat_axis.minor_tick_mark = XL_TICK_MARK.NONE
    cat_axis.format.line.fill.background()
    cat_axis.tick_labels.font.name = "Roboto"
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)

    plot = chart.plots[0]
    series = plot.series[0]
    series.gap_width = 80
    plot.has_data_labels = True
    data_labels = series.data_labels
    data_labels.show_value = True
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = "0.0"
    data_labels.font.name = "Roboto"
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = RGBColor(0, 0, 0)

    dark_blue = RGBColor(21, 57, 96)
    for idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = dark_blue

# Duration chart
def update_duration_chart(slide, df_duration):
    """
    Horizontal bar chart for Duration using Duration!J22:K31.
    - Categories from column J, values from column K (one decimal)
    - Dark blue bars, labels outside end
    """
    start_row =1  # Excel row 22 (0-based)
    end_row = 12    # Excel row 31 (exclusive)
    cat_col = 13     # column J
    val_col = 14    # column K

    categories = []
    values = []
    for cat, val in zip(
        df_duration.iloc[start_row:end_row, cat_col],
        df_duration.iloc[start_row:end_row, val_col],
    ):
        categories.append("" if pd.isna(cat) else str(cat))
        try:
            values.append(None if pd.isna(val) else float(val))
        except Exception:
            values.append(None)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Duration", values)

    left = Inches(0.81)
    top = Inches(2.31)
    width = Inches(5.4)
    height = Inches(5.3)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        left, top, width, height,
        chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = False
    chart.has_title = False

    val_axis = chart.value_axis
    val_axis.major_tick_mark = XL_TICK_MARK.NONE
    val_axis.minor_tick_mark = XL_TICK_MARK.NONE
    val_axis.format.line.fill.background()
    val_axis.tick_labels.number_format_is_linked = False
    val_axis.tick_labels.number_format = ";;;"
    if val_axis.has_major_gridlines:
        val_axis.major_gridlines.format.line.fill.background()

    max_val = max((v for v in values if v is not None), default=0)
    val_axis.minimum_scale = 0
    val_axis.maximum_scale = max_val * 1.1 if max_val > 0 else 1

    cat_axis = chart.category_axis
    cat_axis.major_tick_mark = XL_TICK_MARK.NONE
    cat_axis.minor_tick_mark = XL_TICK_MARK.NONE
    cat_axis.format.line.fill.background()
    cat_axis.tick_labels.font.name = "Roboto"
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.color.rgb = RGBColor(0, 0, 0)

    plot = chart.plots[0]
    series = plot.series[0]
    series.gap_width = 80
    plot.has_data_labels = True
    data_labels = series.data_labels
    data_labels.show_value = True
    data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    data_labels.number_format = "0.0"
    data_labels.font.name = "Roboto"
    data_labels.font.size = Pt(9)
    data_labels.font.color.rgb = RGBColor(0, 0, 0)

    dark_blue = RGBColor(21, 57, 96)
    for point in series.points:
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = dark_blue


# Logic 
def build_presentation(xlsm_path: str, template_path: str) -> io.BytesIO:
    # Load template PPT
    prs = Presentation(template_path)

    # Read Excel sheets (openpyxl reads .xlsm/.xlsx; macros aren’t executed)
    dfs = pd.read_excel(
        xlsm_path,
        sheet_name=[
            "LeasingInfographic", "CompetitiveMarketPosition", "ZipCodes",
            "DrawDemo", "DistanceTravelled", "Frequency", "Duration", "MileageDemo"
        ],
        engine="openpyxl",
    )
    df_leasing = dfs["LeasingInfographic"]
    df_sheet2 = dfs["CompetitiveMarketPosition"]
    df_zipcodes = dfs["ZipCodes"]
    df_drawdemos = dfs["DrawDemo"]
    df_distance = dfs["DistanceTravelled"]
    df_frequency = dfs["Frequency"]
    df_duration = dfs["Duration"]
    df_mileage = dfs["MileageDemo"]
    #chart update
    update_mileage_chart(prs.slides[10], df_mileage)
    update_ethnicities_chart(prs.slides[7], df_drawdemos)
    update_frequency_chart(prs.slides[18], df_frequency)
    update_duration_chart(prs.slides[18],df_duration)

    variable_mapping = {
        "MTZIP1": df_zipcodes.iloc[0,11],
        "MTZIP2": df_zipcodes.iloc[1,11],
        "MTZIP3": df_zipcodes.iloc[2,11],
        "MTZIP4": df_zipcodes.iloc[   3,11],
        "MTZIP5": df_zipcodes.iloc[4,11],
        "MTZIP6": df_zipcodes.iloc[5,11],
        "MTZIP7": df_zipcodes.iloc[6,11],
        "MTZIP8": df_zipcodes.iloc[7,11],
        "MTZIP9": df_zipcodes.iloc[8,11],
        "MTZIP_10": df_zipcodes.iloc[9,11],
        "EXECANALYSIS5": df_leasing.iloc[36,1],
        "VOANALYSIS13": df_leasing.iloc[34,1],
        "MTANALYSIS14": df_leasing.iloc[33,1],
        "VL10": f"{int(round(df_leasing.iloc[0, 0] * 100, 0))}%",
        "VOP08": "{:,.0f}".format(df_leasing.iloc[0, 3]),
        "LD08": f"{int(round(df_leasing.iloc[3, 3] * 100, 0))}%",
        "MT08": f"{int(round(df_leasing.iloc[6, 3] * 100, 0))}%",
        "VF08": df_leasing.iloc[11, 3],
        "HH08": f"{int(round(df_leasing.iloc[0, 7] * 100, 0))}%",
        "HHI08": "${:,.0f}".format(df_leasing.iloc[3, 7]),
        "HHIMSA08": "${:,.0f}".format(df_leasing.iloc[3, 8]),
        "CD08": f"{int(round(df_leasing.iloc[6, 7] * 100, 0))}%",
        "VC08": f"{int(round(df_leasing.iloc[9, 7] * 100, 0))}%",
        "DT08": df_leasing.iloc[11, 7],
        "ZIP1": df_leasing.iloc[3, 0],
        "ZIP2": df_leasing.iloc[4, 0],
        "ZIP3": df_leasing.iloc[5, 0],
        "ZIP4": df_leasing.iloc[6, 0],
        "ZIP5": df_leasing.iloc[7, 0],
        "ZIPANALYSIS15": df_zipcodes.iloc[0, 14],
        "DDANALYSIS12": df_drawdemos.iloc[0, 18],
        "CMPANALYSIS10": df_sheet2.iloc[0, 9],
        "XXXXX": df_leasing.iloc[25, 0],
        "MILANALYSIS11": df_mileage.iloc[0, 3],
        "DURANALYSIS39": df_leasing.iloc[32, 1],
        "FREQANALYSIS38": df_leasing.iloc[33, 1],
        "DTANALYSIS37": df_distance.iloc[0, 12],

        # Mileage Demos mapping
        "MA121": f"{df_mileage.iloc[2, 2]:,.0f}", "MB121": f"{df_mileage.iloc[2, 3]:,.0f}", "MC121": f"{df_mileage.iloc[2, 4]:,.0f}", "MD121": f"{df_mileage.iloc[2, 5]:,.0f}",
        "MA122": f"{df_mileage.iloc[3, 2]:,.0f}", "MB122": f"{df_mileage.iloc[3, 3]:,.0f}", "MC122": f"{df_mileage.iloc[3, 4]:,.0f}", "MD122": f"{df_mileage.iloc[3, 5]:,.0f}",
        "MA123": f"{df_mileage.iloc[4, 2] * 100:.1f}%",  "MB123": f"{df_mileage.iloc[4, 3] * 100:.1f}%",  "MC123": f"{df_mileage.iloc[4, 4] * 100:.1f}%",  "MD123": f"{df_mileage.iloc[4, 5] * 100:.1f}%",
        "MA124": f"{df_mileage.iloc[5, 2] * 100:.1f}%",  "MB124": f"{df_mileage.iloc[5, 3] * 100:.1f}%",  "MC124": f"{df_mileage.iloc[5, 4] * 100:.1f}%",  "MD124": f"{df_mileage.iloc[5, 5] * 100:.1f}%",
        "MA125": f"{df_mileage.iloc[6, 2] * 100:.1f}%",  "MB125": f"{df_mileage.iloc[6, 3] * 100:.1f}%",  "MC125": f"{df_mileage.iloc[6, 4] * 100:.1f}%",  "MD125": f"{df_mileage.iloc[6, 5] * 100:.1f}%",
        "MA126": f"{df_mileage.iloc[7, 2] * 100:.1f}%",  "MB126": f"{df_mileage.iloc[7, 3] * 100:.1f}%",  "MC126": f"{df_mileage.iloc[7, 4] * 100:.1f}%",  "MD126": f"{df_mileage.iloc[7, 5] * 100:.1f}%",
        "MA127": f"{df_mileage.iloc[8, 2] * 100:.1f}%",  "MB127": f"{df_mileage.iloc[8, 3] * 100:.1f}%",  "MC127": f"{df_mileage.iloc[8, 4] * 100:.1f}%",  "MD127": f"{df_mileage.iloc[8, 5] * 100:.1f}%",
        "MA128": f"{df_mileage.iloc[9, 2] * 100:.1f}%",  "MB128": f"{df_mileage.iloc[9, 3] * 100:.1f}%",  "MC128": f"{df_mileage.iloc[9, 4] * 100:.1f}%",  "MD128": f"{df_mileage.iloc[9, 5] * 100:.1f}%",
        "MA129": f"{df_mileage.iloc[10, 2]:.1f}", "MB129": f"{df_mileage.iloc[10, 3]:.1f}", "MC129": f"{df_mileage.iloc[10, 4]:.1f}", "MD129": f"{df_mileage.iloc[10, 5]:.1f}",
        "MA130": f"{df_mileage.iloc[11, 2] * 100:.1f}%", "MB130": f"{df_mileage.iloc[11, 3] * 100:.1f}%", "MC130": f"{df_mileage.iloc[11, 4] * 100:.1f}%", "MD130": f"{df_mileage.iloc[11, 5] * 100:.1f}%",
        "MA131": f"{df_mileage.iloc[12, 2] * 100:.1f}%", "MB131": f"{df_mileage.iloc[12, 3] * 100:.1f}%", "MC131": f"{df_mileage.iloc[12, 4] * 100:.1f}%", "MD131": f"{df_mileage.iloc[12, 5] * 100:.1f}%",
        "MA132": f"{df_mileage.iloc[13, 2] * 100:.1f}%", "MB132": f"{df_mileage.iloc[13, 3] * 100:.1f}%", "MC132": f"{df_mileage.iloc[13, 4] * 100:.1f}%", "MD132": f"{df_mileage.iloc[13, 5] * 100:.1f}%",
        "MA133": f"{df_mileage.iloc[14, 2] * 100:.1f}%", "MB133": f"{df_mileage.iloc[14, 3] * 100:.1f}%", "MC133": f"{df_mileage.iloc[14, 4] * 100:.1f}%", "MD133": f"{df_mileage.iloc[14, 5] * 100:.1f}%",
        "MA134": f"{df_mileage.iloc[15, 2] * 100:.1f}%", "MB134": f"{df_mileage.iloc[15, 3] * 100:.1f}%", "MC134": f"{df_mileage.iloc[15, 4] * 100:.1f}%", "MD134": f"{df_mileage.iloc[15, 5] * 100:.1f}%",
        "MA135": "${:,.0f}".format(df_mileage.iloc[16, 2]), "MB135": "${:,.0f}".format(df_mileage.iloc[16, 3]), "MC135": "${:,.0f}".format(df_mileage.iloc[16, 4]), "MD135": "${:,.0f}".format(df_mileage.iloc[16, 5]),
        "MA136": f"{df_mileage.iloc[17, 2] * 100:.1f}%", "MB136": f"{df_mileage.iloc[17, 3] * 100:.1f}%", "MC136": f"{df_mileage.iloc[17, 4] * 100:.1f}%", "MD136": f"{df_mileage.iloc[17, 5] * 100:.1f}%",
        "MA137": f"{df_mileage.iloc[18, 2] * 100:.1f}%", "MB137": f"{df_mileage.iloc[18, 3] * 100:.1f}%", "MC137": f"{df_mileage.iloc[18, 4] * 100:.1f}%", "MD137": f"{df_mileage.iloc[18, 5] * 100:.1f}%",
        "MA138": f"{df_mileage.iloc[19, 2] * 100:.1f}%", "MB138": f"{df_mileage.iloc[19, 3] * 100:.1f}%", "MC138": f"{df_mileage.iloc[19, 4] * 100:.1f}%", "MD138": f"{df_mileage.iloc[19, 5] * 100:.1f}%",
        "MA139": f"{df_mileage.iloc[20, 2] * 100:.1f}%", "MB139": f"{df_mileage.iloc[20, 3] * 100:.1f}%", "MC139": f"{df_mileage.iloc[20, 4] * 100:.1f}%", "MD139": f"{df_mileage.iloc[20, 5] * 100:.1f}%",
        "MA140": f"{df_mileage.iloc[21, 2] * 100:.1f}%", "MB140": f"{df_mileage.iloc[21, 3] * 100:.1f}%", "MC140": f"{df_mileage.iloc[21, 4] * 100:.1f}%", "MD140": f"{df_mileage.iloc[21, 5] * 100:.1f}%",
        "MA141": f"{df_mileage.iloc[22, 2] * 100:.1f}%", "MB141": f"{df_mileage.iloc[22, 3] * 100:.1f}%", "MC141": f"{df_mileage.iloc[22, 4] * 100:.1f}%", "MD141": f"{df_mileage.iloc[22, 5] * 100:.1f}%",
        "MA142": f"{df_mileage.iloc[23, 2] * 100:.1f}%", "MB142": f"{df_mileage.iloc[23, 3] * 100:.1f}%", "MC142": f"{df_mileage.iloc[23, 4] * 100:.1f}%", "MD142": f"{df_mileage.iloc[23, 5] * 100:.1f}%",
        "MA143": f"{df_mileage.iloc[24, 2] * 100:.1f}%", "MB143": f"{df_mileage.iloc[24, 3] * 100:.1f}%", "MC143": f"{df_mileage.iloc[24, 4] * 100:.1f}%", "MD143": f"{df_mileage.iloc[24, 5] * 100:.1f}%",
        "MA144": f"{df_mileage.iloc[25, 2] * 100:.1f}%", "MB144": f"{df_mileage.iloc[25, 3] * 100:.1f}%", "MC144": f"{df_mileage.iloc[25, 4] * 100:.1f}%", "MD144": f"{df_mileage.iloc[25, 5] * 100:.1f}%",
        "MA145": f"{df_mileage.iloc[26, 2] * 100:.1f}%", "MB145": f"{df_mileage.iloc[26, 3] * 100:.1f}%", "MC145": f"{df_mileage.iloc[26, 4] * 100:.1f}%", "MD145": f"{df_mileage.iloc[26, 5] * 100:.1f}%"
    }

   


    # Replace placeholders across shapes and tables
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in variable_mapping.items():
                            if key in run.text:
                                run.text = run.text.replace(key, str(value))

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for key, value in variable_mapping.items():
                            if key in cell.text:
                                cell.text = cell.text.replace(key, str(value))
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.alignment = PP_ALIGN.CENTER
                                    for run in paragraph.runs:
                                        run.font.name = "Roboto"
                                        run.font.size = Pt(9)
                                        run.font.color.rgb = RGBColor(0, 0, 0)

    # Table copy/format rules
    slides_to_update = {
        9: df_sheet2,
        11: df_drawdemos,
        14: df_zipcodes,
        30: df_distance,
        31: df_frequency,
        32: df_duration
    }

    formatting_rules = {
        "CompetitiveMarketPosition": {
            "percent_columns": [3, 5, 6],
            "currency_columns": [2],
            "thousands_columns": [1]
          },
          "ZipCodes": {
              "percent_columns": [4, 5, 8],
              "currency_columns": [9],
              "thousands_columns": [6, 7]
          },
          "DrawDemo": {
              "percent_rows": [
                "18-24", "25-34", "35-44", "45-54", "55-64", "65+",
                "Less than $50,000", "$50,000-$74,999", "$75,000-$99,999",
                "$100,000-$149,999", "$150,000 or more", "CHILDREN IN HOUSEHOLD",
                "Less than college", "Some college", "College degree", "Post-graduate degree",
                "Caucasian/White", "African-American/Black", "Hispanic/Latino",
                "Asian", "Other"
            ],
              "currency_rows": ["HOUSEHOLD INCOME", "Average HH Income"]
          },
          "DistanceTravelled": {
              "percent_columns": [1, 2, 3, 4, 5],
              "decimal_columns": [6, 7]
          },
          "Frequency": {
              "percent_columns": [1, 2, 3],
              "decimal_columns": [4]
          },
          "Duration": {
              "percent_columns": [1, 2, 3, 4],
              "decimal_columns": [5]
          },
          "CompetitiveMarketPosition": {
              "percent_columns": [3, 5, 6],
              "currency_columns": [2],
              "thousands_columns": [1],
              "decimal_columns": [4]
          }
      }

    def _append_row_clone(tbl):
        """Clone the last row when python-pptx doesn't expose add_row()."""
        if len(tbl.rows) == 0:
            return
        last_row = tbl.rows[len(tbl.rows) - 1]
        clone_tr = deepcopy(last_row._tr)
        tbl._tbl.append(clone_tr)

    for slide_number, df_data in slides_to_update.items():
        slide = prs.slides[slide_number]
        sheet_name = [name for name, df in dfs.items() if df.equals(df_data)][0]
        rules = formatting_rules.get(sheet_name, {})
        table = next((s.table for s in slide.shapes if getattr(s, "has_table", False)), None)
        if not table:
            continue

        rows, cols = df_data.shape
        drawdemo_section_rows = {"AGE", "HOUSEHOLD INCOME", "EDUCATION", "ETHNICITY"}

        # Ensure table has enough rows for incoming data (header + data rows)
        try:
            needed_rows = rows + 1  # include header
            while len(table.rows) < needed_rows:
                try:
                    table.rows.add_row()
                except Exception:
                    _append_row_clone(table)
        except Exception:
            pass  # if cloning unsupported, we'll guard per-row below

        # Header styling where needed
        if sheet_name in ("DrawDemo", "DistanceTravelled", "Frequency", "Duration", "CompetitiveMarketPosition", "ZipCodes"):
            for col_index, col_name in enumerate(df_data.columns):
                if sheet_name == "DrawDemo" and col_index == 0:
                    continue  # keep first column header from template formatting
                if col_index < len(table.columns):
                    cell = table.cell(0, col_index)
                    cell.text = str(col_name)
                    for paragraph in cell.text_frame.paragraphs:
                        if sheet_name == "ZipCodes":
                            paragraph.alignment = PP_ALIGN.CENTER
                        elif sheet_name == "CompetitiveMarketPosition":
                            paragraph.alignment = PP_ALIGN.CENTER
                        else:
                            paragraph.alignment = PP_ALIGN.CENTER if col_index == 0 else PP_ALIGN.LEFT
                        for run in paragraph.runs:
                            run.font.name = "Roboto"
                            run.font.size = Pt(9)
                            # ZipCodes header text should be white/bold; others remain as before
                            if sheet_name == "ZipCodes":
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.bold = True
                                paragraph.alignment = PP_ALIGN.CENTER
                            else:
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.bold = True

        max_data_rows = rows  # we will ensure rows exist as we go
        for row_index in range(max_data_rows):
            # make sure target row exists (header offset of +1)
            try:
                while row_index + 1 >= len(table.rows):
                    try:
                        table.rows.add_row()
                    except Exception:
                        _append_row_clone(table)
            except Exception:
                pass
            # if still not enough rows, stop to avoid index errors
            if row_index + 1 >= len(table.rows):
                break

            for col_index in range(min(cols, len(table.columns))):
                value = df_data.iloc[row_index, col_index]
                value = "" if pd.isna(value) else value
                row_label = str(df_data.iloc[row_index, 0])
                normalized_label = row_label.strip().upper()

                if sheet_name == "DrawDemo" and normalized_label in drawdemo_section_rows:
                    # Leave section header rows untouched (preserve template formatting/height)
                    try:
                        table.rows[row_index + 1].height = Inches(0.17)
                    except Exception:
                        pass
                    continue

                if isinstance(value, numbers.Number):
                    if sheet_name == "DrawDemo":
                        if any(keyword in row_label for keyword in rules.get("percent_rows", [])):
                            formatted_value = f"{round(value * 100, 1)}%"
                        elif any(keyword in row_label for keyword in rules.get("currency_rows", [])):
                            formatted_value = "${:,.0f}".format(value)
                        else:
                            formatted_value = str(value)
                    else:
                        if col_index in rules.get("percent_columns", []):
                            formatted_value = f"{round(value * 100, 1)}%"
                        elif col_index in rules.get("currency_columns", []):
                            formatted_value = "${:,.0f}".format(value)
                        elif col_index in rules.get("thousands_columns", []):
                            formatted_value = "{:,.0f}".format(value)
                        elif col_index in rules.get("decimal_columns", []):
                            formatted_value = f"{value:.1f}"
                        else:
                            formatted_value = str(value)
                else:
                    formatted_value = str(value)

                if sheet_name == "ZipCodes" and row_index == max_data_rows - 1:
                    # Last row formatting for Zip table
                    if col_index == 1:
                        formatted_value = "AVERAGE/TOTALS"
                    if col_index == 2:
                        formatted_value = ""
                    if col_index == 3:
                        try:
                            formatted_value = f"{float(value):.1f}"
                        except Exception:
                            pass
                    try:
                        table.rows[row_index + 1].height = Inches(0.17)
                    except Exception:
                        pass

                if sheet_name == "DrawDemo" and col_index == 0:
                    continue  # preserve first column labels/formatting in template
                cell = table.cell(row_index + 1, col_index)
                # avoid wrapping which can increase row height
                cell.text_frame.word_wrap = False
                cell.text_frame.auto_size = None
                cell.text_frame.margin_top = 0
                cell.text_frame.margin_bottom = 0
                cell.text_frame.margin_left = 0
                cell.text_frame.margin_right = 0
                cell.text = formatted_value
                for paragraph in cell.text_frame.paragraphs:
                    if sheet_name == "ZipCodes":
                        paragraph.alignment = PP_ALIGN.RIGHT if col_index >= 3 else PP_ALIGN.LEFT
                    if sheet_name == "CompetitiveMarketPosition":
                        paragraph.alignment = PP_ALIGN.RIGHT if col_index >= 3 else PP_ALIGN.LEFT
                    paragraph.space_before = Pt(0)
                    paragraph.space_after = Pt(0)
                    paragraph.line_spacing = 1.0
                    for run in paragraph.runs:
                        run.font.name = "Roboto"
                        run.font.size = Pt(9)
                        if sheet_name == "ZipCodes" and row_index == max_data_rows - 1:
                            run.font.bold = True
                            paragraph.space_before = Pt(0)
                            paragraph.space_after = Pt(0)
                            try:
                                paragraph.line_spacing = 1.0
                            except Exception:
                                pass
                        if sheet_name in ("DistanceTravelled", "Frequency", "Duration") and row_index == 0:
                            run.font.color.rgb = RGBColor(255, 255, 255)
                            run.font.bold = True
                        else:
                            run.font.color.rgb = RGBColor(0, 0, 0)
        if sheet_name == "DrawDemo":
            target_height = Inches(0.17)
            for row in table.rows:
                row.height = target_height
        if sheet_name == "ZipCodes":
            target_height = Inches(0.17)
            for row in table.rows:
                row.height = target_height
                try:
                    row._tr.set("h", str(int(target_height)))
                    row._tr.set("hRule", "exact")
                except Exception:
                    pass
            try:
                last_idx = min(max_data_rows, len(table.rows) - 1)
                table.rows[last_idx].height = target_height
                table.rows[last_idx]._tr.set("h", str(int(target_height)))
                table.rows[last_idx]._tr.set("hRule", "exact")
            except Exception:
                pass
        def _apply_banding(table, avg_keywords):
            """Apply alternating light-blue/white rows and gray average row; preserve template header rows."""
            light_blue = RGBColor(0xDD, 0xEB, 0xF7)
            gray = RGBColor(0xD9, 0xD9, 0xD9)
            black = RGBColor(0x00, 0x00, 0x00)
            avg_row_idx = None
            for r_idx in range(1, len(table.rows)):
                label = table.cell(r_idx, 0).text.strip().lower()
                if any(k in label for k in avg_keywords):
                    avg_row_idx = r_idx
                    break
            if avg_row_idx is None:
                avg_row_idx = len(table.rows) - 1

            # Row indices:
            # 0 = header, 1 = sub-header (keep template dark blue), 2 = first data row (force white)
            # start banding from row 3 (index 3) downward.
            first_data_row = 2
            for c_idx in range(len(table.columns)):
                cell = table.cell(first_data_row, c_idx)
                cell.fill.background()
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = black
                        run.font.bold = False

            start_band_row = 3
            for offset, r_idx in enumerate(range(start_band_row, avg_row_idx)):
                for c_idx in range(len(table.columns)):
                    cell = table.cell(r_idx, c_idx)
                    if offset % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = light_blue
                    else:
                        cell.fill.background()
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = black
                            run.font.bold = False

            # average row to gray
            for c_idx in range(len(table.columns)):
                cell = table.cell(avg_row_idx, c_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = gray
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = black

        if sheet_name in ("CompetitiveMarketPosition", "DistanceTravelled", "Frequency", "Duration"):
            _apply_banding(
                table,
                avg_keywords=["average of tested locations", "average/totals"]
            )

    # Write PPT to memory buffer and return
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# Routers
@app.route("/")
def index():
    if not os.path.exists(PPT_TEMPLATE_PATH):
        flash("Template not found: put TT_report.pptx beside app.py")
    return render_template_string(
        HTML,
        template_name=os.path.basename(PPT_TEMPLATE_PATH),
    )

@app.route("/download-template")
def download_template():
    return redirect(SHAREPOINT_TEMPLATE_URL)

# @app.route("/generate", methods=["POST"])
# def generate():
#     file = request.files.get("xlsm")
#     if not file or file.filename == "":
#         flash("Please choose an .xlsm or .xlsx file.")
#         return redirect(url_for("index"))

#     ext = os.path.splitext(file.filename)[1].lower()
#     if ext not in ALLOWED_EXCEL_EXTS:
#         flash("Unsupported file type. Upload .xlsm or .xlsx.")
#         return redirect(url_for("index"))

#     with tempfile.TemporaryDirectory() as tmpdir:
#         safe_name = secure_filename(file.filename)
#         xlsm_path = os.path.join(tmpdir, safe_name)
#         file.save(xlsm_path)

#         try:
#             output = build_presentation(xlsm_path, PPT_TEMPLATE_PATH)
#         except Exception as e:
#             flash(f"Error generating PPT: {e}")
#             return redirect(url_for("index"))

#     ts = datetime.now().strftime("%Y-%m-%d_%H-%M")
#     return send_file(
#         output,
#         as_attachment=True,
#         download_name=f"TT_report_{ts}.pptx",
#         mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
#     )

@app.route("/generate", methods=["POST"])
def generate():
    file = request.files.get("xlsm")
    if not file:
        return "No file uploaded", 400

    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in {".xlsm", ".xlsx"}:
        return "Invalid file type", 400

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            path = os.path.join(tmpdir, secure_filename(file.filename))
            file.save(path)
            output = build_presentation(path, PPT_TEMPLATE_PATH)

    except Exception as e:
        import traceback
        traceback.print_exc()

        return f"ERROR: {e}", 500

    ts = datetime.now().strftime("%Y-%m-%d_%H-%M")
    return send_file(
        output,
        as_attachment=True,
        download_name=f"TT_report_{ts}.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# --------------- Run locally ---------------
if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
